from pptx import Presentation
from pathlib import Path


def create_slide_from_scratch(tmp_dir: Path) -> Path:
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    # title.text = "Hello, World!"
    title.text = "こんにちは, 世界!"
    subtitle.text = "python-pptx was here!"

    p = tmp_dir.joinpath("scratch.pptx")
    prs.save(p)
    return p
